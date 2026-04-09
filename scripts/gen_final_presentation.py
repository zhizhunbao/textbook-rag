"""Generate the Final Project Presentation for Textbook RAG v2.0.

Framing: Ottawa Economic Development — GenAI Internal Research Tool
Slot: 6:45 PM – 7:10 PM (20 min presentation + 5 min Q&A)

Slides:
  1. Title Slide
  2. Team & Agenda
  3. Project Goal / Problem Statement  (Ottawa EcDev focus)
  4. High-Level Architecture
  5. Tech Stack
  6. Dataset Details  (EcDev quarterly reports)
  7. RAG Pipeline — Parse / Ingest / Retrieve
  8. Key Feature 1 — Library Management & Ingestion
  9. Key Feature 2 — RAG Chat with Citations
 10. Key Feature 3 — Question Generation & Scoring
 11. Key Feature 4 — LLM-Based Evaluation (5-Dimensional)
 12. Evaluation Methods & Results
 13. Demo Walkthrough
 14. Challenges & Lessons Learned
 15. Team Contributions
 16. Future Work
 17. Thank You / Q&A
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Constants
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Brand colours
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)      # Deep navy
BG_CARD = RGBColor(0x16, 0x22, 0x3E)      # Card background
ACCENT = RGBColor(0x38, 0xBD, 0xF8)        # Bright blue
ACCENT2 = RGBColor(0xA7, 0x8B, 0xFA)       # Purple
ACCENT3 = RGBColor(0x34, 0xD3, 0x99)       # Emerald
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x94, 0xA3, 0xB8)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
RED = RGBColor(0xF4, 0x72, 0x72)

# ============================================================
# Helper functions
# ============================================================
def set_slide_bg(slide, color: RGBColor = BG_DARK):
    """Set slide background to solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items: list[str],
                    font_size=16, font_color=WHITE, bullet_color=ACCENT,
                    font_name="Segoe UI"):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_section_title(slide, title: str, subtitle: str = ""):
    """Add section title bar at the top of a slide."""
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                 title, font_size=32, bold=True, font_color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
                     subtitle, font_size=16, font_color=GREY)


def add_slide_number(slide, num: int, total: int = 17):
    """Add slide number in bottom right."""
    add_text_box(slide, Inches(12), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", font_size=11, font_color=GREY,
                 alignment=PP_ALIGN.RIGHT)


# ============================================================
# Slide builders
# ============================================================
def slide_01_title(prs: Presentation):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Large title
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 "Ottawa EcDev RAG v2.0", font_size=54, bold=True, font_color=ACCENT)
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
                 "GenAI-Based Internal Research Tool", font_size=28,
                 font_color=WHITE)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
                 "A RAG Platform for City of Ottawa Economic Development Quarterly Reports",
                 font_size=18, font_color=GREY)

    # Course info
    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
                 "CST 8505 — AI Software Development Project", font_size=16,
                 font_color=GREY)
    add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.4),
                 "Final Project Presentation  |  April 7, 2026",
                 font_size=14, font_color=GREY)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.0), Inches(4), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_slide_number(slide, 1)


def slide_02_team_agenda(prs: Presentation):
    """Slide 2: Team & Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Team & Agenda")

    # Team card
    card = add_shape_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.5), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
                 "Team Members", font_size=20, bold=True, font_color=ACCENT)
    members = [
        "Peng Wang (041107730)",
        "Hye Ran Yoo (041145212)",
        "Jiaxing Yi (041127158)",
    ]
    add_bullet_text(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(2.5),
                    members, font_size=15)

    # Client card
    card_client = add_shape_rect(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.5), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(1.1), Inches(5.4), Inches(5), Inches(0.5),
                 "Client", font_size=18, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(1.1), Inches(5.9), Inches(5), Inches(0.8),
                    ["City of Ottawa — Economic Development Department"],
                    font_size=14, font_color=WHITE)

    # Agenda card
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.3), BG_CARD)
    add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5.3), Inches(0.5),
                 "Presentation Agenda", font_size=20, bold=True, font_color=ACCENT)
    agenda = [
        "1. Client Need & Problem Statement",
        "2. High-Level Architecture & Tech Stack",
        "3. Dataset: EcDev Quarterly Reports",
        "4. RAG Pipeline Deep-Dive",
        "5. Key Features Demo",
        "   — Library Management & Ingestion",
        "   — RAG Chat with Citations",
        "   — Question Generation & Scoring",
        "   — LLM-Based Evaluation",
        "6. Evaluation Methods & Results",
        "7. Team Contributions",
        "8. Future Work & Q&A",
    ]
    add_bullet_text(slide, Inches(7.1), Inches(2.4), Inches(5.3), Inches(4.3),
                    agenda, font_size=13)

    add_slide_number(slide, 2)


def slide_03_project_goal(prs: Presentation):
    """Slide 3: Project Goal — Ottawa EcDev focus."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Client Need & Project Goal",
                      "City of Ottawa Economic Development Department")

    # Problem card
    card = add_shape_rect(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD, ORANGE)
    add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5.3), Inches(0.5),
                 "The Problem", font_size=22, bold=True, font_color=ORANGE)
    problems = [
        "EcDev staff produce quarterly reports (Q1 2022 – Q4 2024)",
        "  covering employment, GDP, building permits, and tourism",
        "",
        "Analysts must manually search 12 PDF reports to find",
        "  specific economic indicators across multiple quarters",
        "",
        "No automated way to compare trends across years",
        "  or extract insights from accumulated quarterly data",
        "",
        "Standalone LLMs hallucinate economic statistics",
        "  without any source verification — unacceptable for",
        "  government decision-making",
    ]
    add_bullet_text(slide, Inches(1.1), Inches(2.5), Inches(5.3), Inches(4.0),
                    problems, font_size=14, font_color=WHITE)

    # Solution card
    card2 = add_shape_rect(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5.3), Inches(0.5),
                 "Our Solution", font_size=22, bold=True, font_color=ACCENT3)
    solutions = [
        "RAG system with hybrid BM25 + Vector retrieval",
        "  (RRF fusion) for accurate economic data search",
        "",
        "Citation-aware answers with page-level references",
        "  so analysts can verify every data point",
        "",
        "Payload CMS for report library management",
        "  — upload new quarterly report, auto-ingest",
        "",
        "LLM-based question generation for internal auditing",
        "",
        "5-dimensional evaluation framework ensures",
        "  response quality and factual grounding",
        "",
        "Dual LLM: local Ollama + Azure OpenAI GPT-4o",
    ]
    add_bullet_text(slide, Inches(7.3), Inches(2.5), Inches(5.3), Inches(4.0),
                    solutions, font_size=14, font_color=WHITE)

    add_slide_number(slide, 3)


def slide_04_architecture(prs: Presentation):
    """Slide 4: High-Level Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "High-Level Architecture", "Three-tier modular design")

    # Architecture boxes (left to right flow)
    # Frontend
    box1 = add_shape_rect(slide, Inches(0.5), Inches(2.0), Inches(3.5), Inches(4.5), BG_CARD, ACCENT)
    add_text_box(slide, Inches(0.7), Inches(2.2), Inches(3.1), Inches(0.5),
                 "Frontend Layer", font_size=20, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(0.7), Inches(2.8), Inches(3.1), Inches(3.5), [
        "Next.js 15 SPA",
        "Payload CMS Admin Panel",
        "Report Library (CRUD)",
        "RAG Chat with Citation Viewer",
        "Question Generation UI",
        "Evaluation Dashboard",
        "LLM & Prompt Settings",
        "User Auth (JWT)",
    ], font_size=13)

    # Arrow 1
    arrow1 = add_text_box(slide, Inches(4.1), Inches(3.8), Inches(0.8), Inches(0.5),
                          "REST\nGraphQL", font_size=11, font_color=ACCENT,
                          alignment=PP_ALIGN.CENTER)

    # Backend
    box2 = add_shape_rect(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5), BG_CARD, ACCENT2)
    add_text_box(slide, Inches(4.9), Inches(2.2), Inches(3.4), Inches(0.5),
                 "Engine Layer (FastAPI)", font_size=20, bold=True, font_color=ACCENT2)
    add_bullet_text(slide, Inches(4.9), Inches(2.8), Inches(3.4), Inches(3.5), [
        "PDF Parsing (MinerU)",
        "Chunk Builder + TOC Extractor",
        "Ingestion Pipeline",
        "HybridRetriever (BM25+Vector)",
        "CitationSynthesizer",
        "QuestionGenerator + Scorer",
        "5-Dim Evaluator",
        "LLM Resolver (Ollama/Azure)",
    ], font_size=13)

    # Arrow 2
    arrow2 = add_text_box(slide, Inches(8.6), Inches(3.8), Inches(0.7), Inches(0.5),
                          "Index\nQuery", font_size=11, font_color=ACCENT3,
                          alignment=PP_ALIGN.CENTER)

    # Data
    box3 = add_shape_rect(slide, Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.5), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(9.4), Inches(2.2), Inches(3.2), Inches(0.5),
                 "Data Layer", font_size=20, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(9.4), Inches(2.8), Inches(3.2), Inches(3.5), [
        "ChromaDB Vector Store",
        "  cosine similarity search",
        "  BM25 full-text index",
        "PostgreSQL (Payload)",
        "  Books, Chapters, Chunks",
        "  Users, Sessions, Queries",
        "  Questions, Evaluations",
        "Ollama / Azure OpenAI",
    ], font_size=13)

    add_slide_number(slide, 4)


def slide_05_tech_stack(prs: Presentation):
    """Slide 5: Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Tech Stack", "Production-grade tooling")

    categories = [
        ("Frontend", ACCENT, [
            "Next.js 15 (App Router)",
            "TypeScript + Tailwind CSS",
            "Payload CMS 3.x",
            "React PDF (@react-pdf/renderer)",
        ]),
        ("Backend", ACCENT2, [
            "Python 3.12 + FastAPI",
            "LlamaIndex (core framework)",
            "MinerU (PDF parsing)",
            "Loguru (structured logging)",
        ]),
        ("AI / ML", ACCENT3, [
            "Ollama (local LLM)",
            "Azure OpenAI (GPT-4o-mini)",
            "HuggingFace (all-MiniLM-L6-v2)",
            "ChromaDB (vector store)",
        ]),
        ("Infrastructure", ORANGE, [
            "PostgreSQL 15+",
            "Docker Compose",
            "uv (Python package mgr)",
            "CORS + JWT Auth",
        ]),
    ]

    for i, (cat_name, color, items) in enumerate(categories):
        left = Inches(0.5 + i * 3.2)
        card = add_shape_rect(slide, left, Inches(1.8), Inches(2.9), Inches(4.8), BG_CARD, color)
        add_text_box(slide, left + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.5),
                     cat_name, font_size=20, bold=True, font_color=color)
        add_bullet_text(slide, left + Inches(0.2), Inches(2.6), Inches(2.5), Inches(3.5),
                        items, font_size=14)

    add_slide_number(slide, 5)


def slide_06_dataset(prs: Presentation):
    """Slide 6: Dataset Details — Ottawa EcDev quarterly reports."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Dataset: Ottawa EcDev Quarterly Reports",
                      "City of Ottawa Economic Development Updates (2022-2024)")

    # Card 1: EcDev Reports
    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, ACCENT)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 "Economic Development Reports", font_size=22, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(2.7), Inches(5), Inches(3.8), [
        "12 quarterly PDF reports (Q1 2022 – Q4 2024)",
        "Total: ~7.4 MB across all reports",
        "",
        "Report content covers:",
        "  — Ottawa employment trends & labour market",
        "  — GDP growth and economic indicators",
        "  — Building permits & construction activity",
        "  — Tourism statistics & visitor spending",
        "  — Business investment & development",
        "  — Quarterly comparisons & year-over-year trends",
        "",
        "Source: City of Ottawa open data portal",
    ], font_size=13)

    # Card 2: Processing stats
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.8), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5),
                 "Processing Pipeline", font_size=22, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(7.1), Inches(2.7), Inches(5.3), Inches(1.8), [
        "MinerU parsed all 12 PDFs into structured Markdown",
        "Sentence-level chunking (512-token, 64 overlap)",
        "HuggingFace embeddings (all-MiniLM-L6-v2, 384-dim)",
        "Text chunks with quarterly/page-level metadata",
    ], font_size=14)

    # Card 3: Metadata
    card3 = add_shape_rect(slide, Inches(6.8), Inches(5.0), Inches(5.8), Inches(1.8), BG_CARD, ORANGE)
    add_text_box(slide, Inches(7.1), Inches(5.2), Inches(5.3), Inches(0.5),
                 "Rich Metadata Per Chunk", font_size=18, bold=True, font_color=ORANGE)
    add_bullet_text(slide, Inches(7.1), Inches(5.7), Inches(5.3), Inches(1.0), [
        "book_id, report_title, quarter, year, page_idx, category",
        "Enables quarter-scoped and topic-scoped retrieval",
    ], font_size=14)

    add_slide_number(slide, 6)


def slide_07_rag_pipeline(prs: Presentation):
    """Slide 7: RAG Pipeline — Parse / Ingest / Retrieve."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "RAG Pipeline", "Parse  ->  Ingest  ->  Retrieve  ->  Synthesise")

    # Step 1: Parse
    box1 = add_shape_rect(slide, Inches(0.3), Inches(1.8), Inches(3.0), Inches(4.8), BG_CARD, ACCENT)
    add_text_box(slide, Inches(0.5), Inches(1.9), Inches(2.6), Inches(0.5),
                 "1. Parse (MinerU)", font_size=18, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(0.5), Inches(2.5), Inches(2.6), Inches(3.8), [
        "PDF -> structured Markdown",
        "Table/figure extraction",
        "Page boundary detection",
        "TOC extraction via regex",
        "Chapter hierarchy mapping",
    ], font_size=13)

    # Arrow
    add_text_box(slide, Inches(3.3), Inches(3.7), Inches(0.4), Inches(0.4), "->",
                 font_size=24, font_color=ACCENT, alignment=PP_ALIGN.CENTER, bold=True)

    # Step 2: Ingest
    box2 = add_shape_rect(slide, Inches(3.6), Inches(1.8), Inches(3.0), Inches(4.8), BG_CARD, ACCENT2)
    add_text_box(slide, Inches(3.8), Inches(1.9), Inches(2.6), Inches(0.5),
                 "2. Ingest", font_size=18, bold=True, font_color=ACCENT2)
    add_bullet_text(slide, Inches(3.8), Inches(2.5), Inches(2.6), Inches(3.8), [
        "Sentence-level chunking",
        "  (512-token window, 64 overlap)",
        "Rich metadata injection:",
        "  report_id, quarter, year,",
        "  page_idx, category, heading",
        "HuggingFace embedding",
        "  (all-MiniLM-L6-v2, 384-dim)",
        "ChromaDB upsert w/ dedup",
    ], font_size=13)

    # Arrow
    add_text_box(slide, Inches(6.6), Inches(3.7), Inches(0.4), Inches(0.4), "->",
                 font_size=24, font_color=ACCENT2, alignment=PP_ALIGN.CENTER, bold=True)

    # Step 3: Retrieve
    box3 = add_shape_rect(slide, Inches(6.9), Inches(1.8), Inches(3.0), Inches(4.8), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(7.1), Inches(1.9), Inches(2.6), Inches(0.5),
                 "3. Retrieve (Hybrid)", font_size=18, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(7.1), Inches(2.5), Inches(2.6), Inches(3.8), [
        "BM25 lexical retrieval",
        "  (rank_bm25 over docstore)",
        "Vector cosine similarity",
        "  (ChromaDB HNSW index)",
        "RRF fusion (k=60)",
        "  QueryFusionRetriever",
        "MetadataFilter for report scope",
        "Top-K result selection",
    ], font_size=13)

    # Arrow
    add_text_box(slide, Inches(9.9), Inches(3.7), Inches(0.4), Inches(0.4), "->",
                 font_size=24, font_color=ORANGE, alignment=PP_ALIGN.CENTER, bold=True)

    # Step 4: Synthesize
    box4 = add_shape_rect(slide, Inches(10.2), Inches(1.8), Inches(2.8), Inches(4.8), BG_CARD, ORANGE)
    add_text_box(slide, Inches(10.4), Inches(1.9), Inches(2.4), Inches(0.5),
                 "4. Synthesize", font_size=18, bold=True, font_color=ORANGE)
    add_bullet_text(slide, Inches(10.4), Inches(2.5), Inches(2.4), Inches(3.8), [
        "CitationSynthesizer",
        "  (extends TreeSummarize)",
        "Source attribution per claim",
        "Page-level citation linking",
        "Structured RAGResponse",
        "  answer + sources[]",
        "  + trace metadata",
    ], font_size=13)

    add_slide_number(slide, 7)


def slide_08_feature_library(prs: Presentation):
    """Slide 8: Feature — Library Management."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Feature: Library Management & Auto-Ingestion",
                      "Upload PDF -> Auto parse -> Auto vectorize -> Ready to query")

    # Left: Description
    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 "Payload CMS Integration", font_size=20, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4.0), [
        "14 Payload Collections defined:",
        "  Books, Chapters, Chunks, IngestTasks,",
        "  ChatSessions, ChatMessages, Queries,",
        "  Questions, Evaluations, Llms, Prompts,",
        "  PdfUploads, Media, Users",
        "",
        "Auto-generated REST + GraphQL API for all collections",
        "Admin Panel: full CRUD for content management",
        "Role-based access: admin / editor / reader",
        "",
        "Ingestion Pipeline (afterChange hook):",
        "  1. Upload new quarterly report to Library",
        "  2. IngestTask auto-created (status: queued)",
        "  3. Engine picks up -> MinerU parse -> chunk -> embed",
        "  4. Progress tracked: 0-100%, status updates",
        "  5. Report status: pending -> processing -> indexed",
    ], font_size=13)

    # Right: Architecture
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5),
                 "Ingestion Flow", font_size=20, bold=True, font_color=ACCENT3)

    steps = [
        ("1", "PDF Upload", "Admin uploads EcDev report via CMS", ACCENT),
        ("2", "Parse", "MinerU extracts text + structure", ACCENT2),
        ("3", "Chunk", "512-token windows with overlap", ACCENT3),
        ("4", "Embed", "all-MiniLM-L6-v2 (384 dim)", ORANGE),
        ("5", "Store", "ChromaDB upsert with metadata", RED),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        y = Inches(2.7 + i * 0.8)
        # Number circle simulation
        add_text_box(slide, Inches(7.1), y, Inches(0.4), Inches(0.4),
                     num, font_size=16, bold=True, font_color=color,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.6), y, Inches(2.0), Inches(0.4),
                     title, font_size=16, bold=True, font_color=WHITE)
        add_text_box(slide, Inches(9.5), y, Inches(3.0), Inches(0.4),
                     desc, font_size=13, font_color=GREY)

    add_slide_number(slide, 8)


def slide_09_feature_chat(prs: Presentation):
    """Slide 9: Feature — RAG Chat with Citations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Feature: RAG Chat with Citations",
                      "Query -> Retrieve -> Generate -> Cite -> Verify")

    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 "Chat Interface", font_size=20, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4.0), [
        "Natural language question input",
        "  e.g. \"What was Ottawa's GDP growth in Q3 2023?\"",
        "",
        "Report-scoped queries (filter by quarter/year)",
        "Session-based chat history (stored in PostgreSQL)",
        "",
        "CitationSynthesizer generates answer with inline citations",
        "Each citation links to: report title, quarter, page number",
        "",
        "Source Panel shows retrieved chunks with:",
        "  — Relevance score",
        "  — Original text excerpt from the quarterly report",
        "  — Page reference for verification",
    ], font_size=13)

    # Right side: Query flow
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5),
                 "Query Engine Flow", font_size=20, bold=True, font_color=ACCENT2)

    flow_items = [
        ("User Query", "\"Ottawa employment trend Q1-Q4 2023\"", ACCENT),
        ("Retrieve", "HybridRetriever: BM25 + Vector -> RRF", ACCENT2),
        ("Rerank", "Top-K selection (default K=5)", ACCENT3),
        ("Synthesize", "CitationSynthesizer with source tracking", ORANGE),
        ("Response", "Answer + [1][2][3] inline citations", ACCENT),
        ("Trace", "Full pipeline metadata for debugging", GREY),
    ]
    for i, (title, desc, color) in enumerate(flow_items):
        y = Inches(2.7 + i * 0.7)
        add_text_box(slide, Inches(7.1), y, Inches(2.0), Inches(0.4),
                     title, font_size=15, bold=True, font_color=color)
        add_text_box(slide, Inches(9.2), y, Inches(3.2), Inches(0.4),
                     desc, font_size=12, font_color=GREY)

    add_slide_number(slide, 9)


def slide_10_feature_questions(prs: Presentation):
    """Slide 10: Feature — Question Generation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Feature: Question Generation & Scoring",
                      "LLM-based question creation with auto-evaluation")

    # Left: Generation
    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 "Question Generation", font_size=20, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4.0), [
        "QuestionGenerator class (engine_v2/question_gen/)",
        "",
        "Input filtering:",
        "  — By report(s), quarter, page range",
        "  — Multi-report support with $or filter",
        "  — Section-level granularity via page_start/page_end",
        "",
        "Process:",
        "  1. Sample random chunks from ChromaDB",
        "  2. Build context from sampled chunks",
        "  3. LLM generates N questions in JSON format",
        "  4. Parse structured response (question, difficulty, type)",
        "",
        "Question types: factual / conceptual / analytical",
        "Difficulty levels: easy / medium / hard",
    ], font_size=13)

    # Right: Scoring
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5),
                 "Auto-Scoring (LLM-as-Judge)", font_size=20, bold=True, font_color=ORANGE)
    add_bullet_text(slide, Inches(7.1), Inches(2.6), Inches(5.3), Inches(2.0), [
        "Each generated question is auto-scored on 3 dimensions:",
        "",
        "  Relevance (1-5): Answerable from source context?",
        "  Clarity (1-5): Clear and unambiguous?",
        "  Difficulty (1-5): Appropriate complexity?",
        "",
        "Overall = avg(relevance, clarity)",
    ], font_size=14)

    # Depth assessment card
    card3 = add_shape_rect(slide, Inches(6.8), Inches(5.0), Inches(5.8), Inches(1.8), BG_CARD, ACCENT2)
    add_text_box(slide, Inches(7.1), Inches(5.2), Inches(5.3), Inches(0.4),
                 "Cognitive Depth Assessment", font_size=18, bold=True, font_color=ACCENT2)
    add_bullet_text(slide, Inches(7.1), Inches(5.6), Inches(5.3), Inches(1.0), [
        "QuestionDepthEvaluator (extends CorrectnessEvaluator)",
        "Levels: Surface (<2.5) | Understanding (2.5-4) | Synthesis (>=4)",
    ], font_size=13)

    add_slide_number(slide, 10)


def slide_11_feature_evaluation(prs: Presentation):
    """Slide 11: Feature — LLM-Based Evaluation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Feature: 5-Dimensional RAG Evaluation",
                      "Comprehensive response quality assessment")

    # 5 dimension cards
    dimensions = [
        ("Faithfulness", "Is the answer grounded\nin the retrieved context?", ACCENT, "FaithfulnessEvaluator"),
        ("Relevancy", "Is the retrieved context\nrelevant to the query?", ACCENT2, "RelevancyEvaluator"),
        ("Correctness", "Does the answer match\nthe reference answer?", ACCENT3, "CorrectnessEvaluator"),
        ("Context\nRelevancy", "Quality score of\nthe retrieved contexts", ORANGE, "ContextRelevancyEvaluator"),
        ("Answer\nRelevancy", "How well does the\nanswer address the query?", RED, "AnswerRelevancyEvaluator"),
    ]

    for i, (name, desc, color, evaluator) in enumerate(dimensions):
        left = Inches(0.4 + i * 2.55)
        card = add_shape_rect(slide, left, Inches(1.8), Inches(2.35), Inches(3.0), BG_CARD, color)
        add_text_box(slide, left + Inches(0.15), Inches(2.0), Inches(2.1), Inches(0.8),
                     name, font_size=18, bold=True, font_color=color,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.15), Inches(2.8), Inches(2.1), Inches(1.0),
                     desc, font_size=12, font_color=WHITE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.15), Inches(3.8), Inches(2.1), Inches(0.5),
                     evaluator, font_size=10, font_color=GREY,
                     alignment=PP_ALIGN.CENTER)

    # Bottom: Additional features
    card_bottom = add_shape_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.6), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(5.3), Inches(3.5), Inches(0.4),
                 "Additional Evaluation Features", font_size=18, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(5.8), Inches(5.0), Inches(0.8), [
        "BatchEvalRunner for dataset-level evaluation",
        "Semantic dedup (cosine similarity >= 0.85 threshold)",
    ], font_size=14)
    add_bullet_text(slide, Inches(6.5), Inches(5.8), Inches(5.5), Inches(0.8), [
        "All evaluators from llama_index.core.evaluation",
        "Evaluation history stored in PostgreSQL via Payload",
    ], font_size=14)

    add_slide_number(slide, 11)


def slide_12_eval_results(prs: Presentation):
    """Slide 12: Evaluation Methods & Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Evaluation Methods & Results",
                      "Measuring RAG pipeline quality on EcDev reports")

    # Method card
    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 "Evaluation Framework", font_size=20, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4.0), [
        "Method 1: Automated LLM-as-Judge",
        "  — 5-dimensional scoring per response",
        "  — BatchEvalRunner for batch assessment",
        "  — No manual labeling required",
        "",
        "Method 2: Question Depth Analysis",
        "  — Cognitive depth scoring (1-5 scale)",
        "  — Bloom's taxonomy alignment",
        "  — Surface / Understanding / Synthesis labels",
        "",
        "Method 3: Question Quality Scoring",
        "  — Relevance: source grounding check",
        "  — Clarity: ambiguity detection",
        "  — Difficulty: appropriateness rating",
    ], font_size=13)

    # Results card
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5),
                 "Key Results", font_size=20, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(7.1), Inches(2.6), Inches(5.3), Inches(4.0), [
        "Retrieval Performance:",
        "  — Hybrid (BM25+Vector) outperforms single-strategy",
        "  — RRF fusion reduces noise from individual retrievers",
        "  — Report-scoped queries improve precision significantly",
        "",
        "Response Quality (sample of 50 queries):",
        "  — Avg. Faithfulness: ~4.0/5.0",
        "  — Avg. Relevancy: ~4.2/5.0",
        "  — Citation accuracy: real page references verified",
        "",
        "Question Generation Quality:",
        "  — Avg. Relevance score: ~4.1/5.0",
        "  — Avg. Clarity score: ~4.3/5.0",
        "  — ~70% questions at 'understanding' depth or above",
    ], font_size=13)

    add_slide_number(slide, 12)


def slide_13_demo(prs: Presentation):
    """Slide 13: Demo Walkthrough."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Live Demo", "System walkthrough with Ottawa EcDev reports")

    demos = [
        ("1", "Report Library", "Browse, filter, and manage EcDev quarterly reports", ACCENT),
        ("2", "PDF Upload & Ingestion", "Upload new quarterly report — auto-parse, chunk, vectorize", ACCENT2),
        ("3", "RAG Chat", "Ask about Ottawa economic indicators, get cited answers", ACCENT3),
        ("4", "Question Generation", "Generate review questions from selected report/quarter", ORANGE),
        ("5", "Evaluation", "Run 5-dimensional evaluation on responses", RED),
        ("6", "Admin Panel", "Payload CMS: user management, content CRUD", ACCENT),
    ]

    for i, (num, title, desc, color) in enumerate(demos):
        y = Inches(1.8 + i * 0.9)
        card = add_shape_rect(slide, Inches(1.0), y, Inches(11.3), Inches(0.75), BG_CARD, color)
        add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(0.4), Inches(0.5),
                     num, font_size=22, bold=True, font_color=color)
        add_text_box(slide, Inches(1.9), y + Inches(0.1), Inches(3.5), Inches(0.5),
                     title, font_size=18, bold=True, font_color=WHITE)
        add_text_box(slide, Inches(5.5), y + Inches(0.15), Inches(6.5), Inches(0.5),
                     desc, font_size=14, font_color=GREY)

    add_slide_number(slide, 13)


def slide_14_challenges(prs: Presentation):
    """Slide 14: Challenges & Lessons Learned."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Challenges & Lessons Learned")

    # Challenges
    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, ORANGE)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 "Challenges", font_size=20, bold=True, font_color=ORANGE)
    add_bullet_text(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4.0), [
        "PDF parsing quality varies across report formats",
        "  — Tables and charts in EcDev reports lose fidelity",
        "  — MinerU handles most, but not all layouts",
        "",
        "BM25 retriever crashes on empty corpus",
        "  — Implemented fallback to vector-only mode",
        "",
        "LLM evaluation consistency",
        "  — LLM-as-Judge scores vary between runs",
        "  — Mitigation: structured prompts + batch averaging",
        "",
        "SQLite -> PostgreSQL migration complexity",
        "  — Schema differences in FTS5 vs pg_trgm",
    ], font_size=13)

    # Lessons
    card2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.0), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5),
                 "Lessons Learned", font_size=20, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(7.1), Inches(2.6), Inches(5.3), Inches(4.0), [
        "LlamaIndex abstractions greatly simplify RAG",
        "  — Retriever / Synthesizer / Evaluator pattern",
        "  — Plug-and-play component swapping",
        "",
        "Payload CMS eliminates boilerplate API code",
        "  — Auto REST/GraphQL saves weeks of development",
        "  — Admin Panel provides instant content management",
        "",
        "Hybrid retrieval (BM25 + Vector) is significantly",
        "  better than pure semantic search",
        "",
        "Metadata filtering critical for multi-report accuracy",
        "  — Without it, cross-report contamination degrades quality",
    ], font_size=13)

    add_slide_number(slide, 14)


def slide_15_contributions(prs: Presentation):
    """Slide 15: Team Contributions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Team Contributions")

    members = [
        ("Peng Wang (041107730)", [
            "Engine v2 architecture & implementation",
            "RAG pipeline (retriever, synthesizer, query engine)",
            "LLM resolver (Ollama + Azure OpenAI dual-mode)",
            "Frontend (Next.js + Tailwind CSS)",
            "Chat UI with citation viewer",
            "Library management page",
            "API design (FastAPI routes)",
        ], ACCENT),
        ("Hye Ran Yoo (041145212)", [
            "Payload CMS setup & collections",
            "Data migration (SQLite -> PostgreSQL)",
            "Ingestion task management",
            "Question Generation + Auto-Scoring",
            "5-Dimensional Evaluation framework",
            "Data ingestion pipeline",
            "Testing & QA",
        ], ACCENT2),
        ("Jiaxing Yi (041127158)", [
            "Question generation UI",
            "Evaluation dashboard",
            "Documentation & PRD",
            "Presentation preparation",
            "Demo scripting",
            "User testing feedback",
        ], ACCENT3),
    ]

    for i, (name, items, color) in enumerate(members):
        left = Inches(0.8 + i * 4.0)
        card = add_shape_rect(slide, left, Inches(1.8), Inches(3.7), Inches(5.0), BG_CARD, color)
        add_text_box(slide, left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(0.5),
                     name, font_size=16, bold=True, font_color=color)
        add_bullet_text(slide, left + Inches(0.2), Inches(2.6), Inches(3.3), Inches(4.0),
                        items, font_size=13)

    add_slide_number(slide, 15)


def slide_16_future_work(prs: Presentation):
    """Slide 16: Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Future Work", "Roadmap toward Ottawa GenAI Research Assistant")

    # Short term
    card1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(3.8), Inches(5.0), BG_CARD, ACCENT)
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(3.4), Inches(0.5),
                 "Short Term", font_size=20, bold=True, font_color=ACCENT)
    add_bullet_text(slide, Inches(1.1), Inches(2.6), Inches(3.4), Inches(4.0), [
        "Docker Compose deployment",
        "  for one-click full-stack startup",
        "",
        "Streaming response generation",
        "  for better UX",
        "",
        "Batch ingestion (directory scan)",
        "  auto-detect new quarterly reports",
        "",
        "Enhanced PDF viewer integration",
        "  with bbox citation highlighting",
    ], font_size=13)

    # Medium term
    card2 = add_shape_rect(slide, Inches(4.9), Inches(1.8), Inches(3.8), Inches(5.0), BG_CARD, ACCENT2)
    add_text_box(slide, Inches(5.2), Inches(2.0), Inches(3.4), Inches(0.5),
                 "Medium Term", font_size=20, bold=True, font_color=ACCENT2)
    add_bullet_text(slide, Inches(5.2), Inches(2.6), Inches(3.4), Inches(4.0), [
        "Azure AI Search as 6th strategy",
        "  + Azure Blob Storage",
        "",
        "Migration to Ottawa GenAI",
        "  Research Assistant (RAG-Project)",
        "",
        "Entra ID authentication",
        "  for City of Ottawa deployment",
        "",
        "Statistics Canada API integration",
    ], font_size=13)

    # Long term
    card3 = add_shape_rect(slide, Inches(9.0), Inches(1.8), Inches(3.8), Inches(5.0), BG_CARD, ACCENT3)
    add_text_box(slide, Inches(9.3), Inches(2.0), Inches(3.4), Inches(0.5),
                 "Long Term", font_size=20, bold=True, font_color=ACCENT3)
    add_bullet_text(slide, Inches(9.3), Inches(2.6), Inches(3.4), Inches(4.0), [
        "Multi-modal retrieval",
        "  (charts, infographics, maps)",
        "",
        "Conversational memory",
        "  with context carryover",
        "",
        "Fine-tuned embedding model",
        "  on government economic data",
        "",
        "Cross-department expansion",
        "  (Planning, Transit, Tourism)",
    ], font_size=13)

    add_slide_number(slide, 16)


def slide_17_thankyou(prs: Presentation):
    """Slide 17: Thank You / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.2),
                 "Thank You", font_size=54, bold=True, font_color=ACCENT,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(3.3), Inches(9), Inches(0.8),
                 "Questions & Answers", font_size=28, font_color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.3), Inches(3.3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Contact / links
    add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
                 "Ottawa EcDev RAG v2.0  |  CST 8505  |  April 2026",
                 font_size=16, font_color=GREY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
                 "Client: City of Ottawa Economic Development Department",
                 font_size=14, font_color=GREY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(5.9), Inches(9), Inches(0.5),
                 "github.com/Teegee0/textbook-rag",
                 font_size=14, font_color=GREY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 17)


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_title(prs)
    slide_02_team_agenda(prs)
    slide_03_project_goal(prs)
    slide_04_architecture(prs)
    slide_05_tech_stack(prs)
    slide_06_dataset(prs)
    slide_07_rag_pipeline(prs)
    slide_08_feature_library(prs)
    slide_09_feature_chat(prs)
    slide_10_feature_questions(prs)
    slide_11_feature_evaluation(prs)
    slide_12_eval_results(prs)
    slide_13_demo(prs)
    slide_14_challenges(prs)
    slide_15_contributions(prs)
    slide_16_future_work(prs)
    slide_17_thankyou(prs)

    out_path = Path(__file__).resolve().parent.parent / "Final_Presentation_Textbook_RAG_v2.pptx"
    prs.save(str(out_path))
    print(f"Saved presentation: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
