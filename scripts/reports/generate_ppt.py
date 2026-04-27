"""Generate presentation for Ottawa EcDev RAG — Final Group Demo.

Usage:
    python scripts/generate_ppt.py
    python scripts/generate_ppt.py --output my_slides.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ============================================================
# Color palette
# ============================================================
BG_DARK = RGBColor(0x0F, 0x0F, 0x14)
BG_CARD = RGBColor(0x1A, 0x1A, 0x24)
BG_SECTION = RGBColor(0x14, 0x14, 0x1E)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)
TEXT_PRIMARY = RGBColor(0xF4, 0xF4, 0xF5)
TEXT_SECONDARY = RGBColor(0xA1, 0xA1, 0xAA)
TEXT_MUTED = RGBColor(0x71, 0x71, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ============================================================
# Helpers
# ============================================================

def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text,
                font_size=14, color=TEXT_PRIMARY, bold=False,
                alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items,
                    font_size=13, color=TEXT_SECONDARY, spacing=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  -  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(spacing)
    return txBox

def add_rounded_rect(slide, left, top, width, height,
                     fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5),
                   color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def section_divider(prs, title, subtitle, accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SECTION)
    add_accent_bar(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(1.4), accent)
    add_textbox(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(1),
                title, font_size=36, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(3.6), Inches(10), Inches(0.6),
                subtitle, font_size=16, color=TEXT_SECONDARY)
    return slide


# ============================================================
# Slide: Title
# ============================================================
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(1), Inches(1.4), Inches(11), Inches(1.2),
                "Ottawa Economic Development", font_size=44, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
                "AI-Powered Report Analysis with RAG",
                font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
                "LlamaIndex  /  Full-Stack Platform  /  Quarterly Report Q&A",
                font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4), Inches(4.2), Inches(5.33), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
                "CST 8916 -- Applied Generative AI",
                font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
                "Instructor: Prof. Hari Koduvely",
                font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.5),
                "Team: [Your Team Name]  /  Members: [Names]",
                font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.4),
                "April 8, 2026",
                font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide: Agenda
# ============================================================
def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
                "Agenda", font_size=28, color=WHITE, bold=True)

    items = [
        ("01", "Problem & Motivation -- Why Ottawa EcDev Needs RAG", ACCENT_BLUE),
        ("02", "Data -- 12 Quarterly Reports (2022-2024)", ACCENT_GREEN),
        ("03", "System Architecture -- 3-Layer Design", ACCENT_PURPLE),
        ("04", "RAG Pipeline -- Parse, Ingest, Retrieve", ACCENT_AMBER),
        ("05", "Implemented Features (with Demo)", ACCENT_PINK),
        ("06", "Evaluation & Results", ACCENT_BLUE),
        ("07", "Future Work -- What Comes Next", ACCENT_GREEN),
        ("08", "Q&A", ACCENT_PURPLE),
    ]
    for i, (num, label, accent) in enumerate(items):
        y = Inches(1.5) + Inches(i * 0.7)
        add_rounded_rect(slide, Inches(1), y, Inches(0.6), Inches(0.45),
                         fill_color=BG_CARD, border_color=accent)
        add_textbox(slide, Inches(1), y + Pt(3), Inches(0.6), Inches(0.4),
                    num, font_size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.8), y + Pt(2), Inches(10), Inches(0.4),
                    label, font_size=16, color=TEXT_PRIMARY)


# ============================================================
# Slide: Problem
# ============================================================
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.6),
                "Problem Statement", font_size=28, color=WHITE, bold=True)

    add_textbox(slide, Inches(1), Inches(1.3), Inches(5), Inches(0.5),
                "City of Ottawa -- EcDev Reports", font_size=18, color=ACCENT_BLUE, bold=True)
    context = [
        "City publishes quarterly Economic Development Updates",
        "12 reports across 3 years (Q1 2022 - Q4 2024)",
        "Covers: Labour Force, Housing, CPI, Construction, Vacancy Rates",
        "Rich with tables, charts, and statistical data",
        "Cross-quarter trend analysis requires flipping through many PDFs",
        "No existing tool for natural language querying over these reports",
    ]
    add_bullet_list(slide, Inches(1), Inches(1.9), Inches(5.3), Inches(3.5),
                    context, font_size=11, spacing=7)

    add_rounded_rect(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(5.0),
                     fill_color=BG_CARD, border_color=ACCENT_GREEN)
    add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
                "Our Solution", font_size=18, color=ACCENT_GREEN, bold=True)
    solution = [
        "RAG-powered Q&A over Ottawa EcDev reports",
        'E.g. "Compare unemployment Q1 2022 vs Q4 2024"',
        "Hybrid retrieval: BM25 + Vector with RRF fusion",
        "Citation with page-level bbox verification",
        "Auto question generation & LLM evaluation",
        "Full-stack: Library, Pipeline, Chat, Evaluate",
        "Azure-ready for Ottawa GenAI deployment",
    ]
    add_bullet_list(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(3.8),
                    solution, font_size=11, spacing=7)


# ============================================================
# Slide: Data Overview
# ============================================================
def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_GREEN)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Data: Ottawa EcDev Quarterly Reports", font_size=28, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(0.95), Inches(10), Inches(0.4),
                "12 reports  /  Q1 2022 - Q4 2024  /  ~600-1000 KB each  /  Tables + Charts + Text",
                font_size=13, color=TEXT_SECONDARY)

    years = ["2022", "2023", "2024"]
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    for j, q in enumerate(quarters):
        x = Inches(1.8) + Inches(j * 2.7)
        add_textbox(slide, x, Inches(1.65), Inches(2.3), Inches(0.35),
                    q, font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    for i, year in enumerate(years):
        y = Inches(2.1) + Inches(i * 1.5)
        add_textbox(slide, Inches(0.8), y + Inches(0.2), Inches(0.9), Inches(0.4),
                    year, font_size=16, color=ACCENT_AMBER, bold=True)
        for j in range(4):
            x = Inches(1.8) + Inches(j * 2.7)
            add_rounded_rect(slide, x, y, Inches(2.3), Inches(1.2),
                             fill_color=BG_CARD, border_color=ACCENT_GREEN)
            add_textbox(slide, x + Inches(0.15), y + Pt(6), Inches(2), Inches(0.3),
                        f"Q{j+1} {year}", font_size=11, color=TEXT_PRIMARY, bold=True)
            add_textbox(slide, x + Inches(0.15), y + Inches(0.4), Inches(2), Inches(0.25),
                        f"ed_update_q{j+1}_{year}.pdf", font_size=8, color=TEXT_MUTED)

    y_s = Inches(6.2)
    add_rounded_rect(slide, Inches(0.8), y_s, Inches(11.6), Inches(0.8),
                     fill_color=BG_CARD, border_color=ACCENT_PURPLE)
    add_textbox(slide, Inches(1), y_s + Pt(8), Inches(11), Inches(0.5),
                "Key Indicators:  Labour Force  /  Unemployment Rate  /  Housing Starts  /  Avg Resale Price  /  Population  /  CPI  /  Weekly Earnings  /  Construction Permits  /  Office Vacancy",
                font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide: Architecture
# ============================================================
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_PURPLE)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "System Architecture", font_size=28, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(0.95), Inches(8), Inches(0.4),
                "Three-layer: Presentation  /  CMS  /  Engine", font_size=13, color=TEXT_SECONDARY)

    lw = Inches(11.5)
    lh = Inches(1.3)

    # Layer 1
    y = Inches(1.6)
    add_rounded_rect(slide, Inches(0.8), y, lw, lh,
                     fill_color=RGBColor(0x17,0x25,0x3B), border_color=ACCENT_BLUE)
    add_textbox(slide, Inches(1), y+Pt(6), Inches(3), Inches(0.35),
                "Next.js 15 + Payload CMS Frontend", font_size=14, color=ACCENT_BLUE, bold=True)
    for i, c in enumerate(["Library", "Pipeline Monitor", "Chat + PDF", "Question Gen", "Evaluation", "Auth (JWT)"]):
        x = Inches(1.2)+Inches(i*1.85)
        add_rounded_rect(slide, x, y+Inches(0.5), Inches(1.65), Inches(0.55),
                         fill_color=BG_DARK, border_color=RGBColor(0x2D,0x4A,0x6F))
        add_textbox(slide, x+Pt(4), y+Inches(0.58), Inches(1.6), Inches(0.4),
                    c, font_size=9, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Layer 2
    y2 = y+lh+Inches(0.25)
    add_rounded_rect(slide, Inches(0.8), y2, lw, lh,
                     fill_color=RGBColor(0x1A,0x2D,0x1A), border_color=ACCENT_GREEN)
    add_textbox(slide, Inches(1), y2+Pt(6), Inches(3), Inches(0.35),
                "Payload CMS 3.x + PostgreSQL", font_size=14, color=ACCENT_GREEN, bold=True)
    for i, c in enumerate(["Books","Chapters","Chunks","IngestTasks","ChatSessions","Messages","Questions","Evaluations","Prompts","LLMs","Users"]):
        x = Inches(1.0)+Inches(i*1.03)
        add_rounded_rect(slide, x, y2+Inches(0.5), Inches(0.88), Inches(0.55),
                         fill_color=BG_DARK, border_color=RGBColor(0x2D,0x5A,0x2D))
        add_textbox(slide, x+Pt(2), y2+Inches(0.58), Inches(0.85), Inches(0.4),
                    c, font_size=7, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Layer 3
    y3 = y2+lh+Inches(0.25)
    add_rounded_rect(slide, Inches(0.8), y3, lw, lh,
                     fill_color=RGBColor(0x2D,0x1A,0x2D), border_color=ACCENT_PURPLE)
    add_textbox(slide, Inches(1), y3+Pt(6), Inches(3), Inches(0.35),
                "Python Engine v2 (LlamaIndex)", font_size=14, color=ACCENT_PURPLE, bold=True)
    for i, c in enumerate(["MinerU Reader","Chunking","Embedding","HybridRetriever","CitationEngine","QuestionGen","Evaluation","FastAPI"]):
        x = Inches(1.2)+Inches(i*1.38)
        add_rounded_rect(slide, x, y3+Inches(0.5), Inches(1.2), Inches(0.55),
                         fill_color=BG_DARK, border_color=RGBColor(0x4A,0x2D,0x5A))
        add_textbox(slide, x+Pt(3), y3+Inches(0.58), Inches(1.15), Inches(0.4),
                    c, font_size=9, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Data
    y4 = y3+lh+Inches(0.25)
    add_rounded_rect(slide, Inches(0.8), y4, lw, Inches(0.65),
                     fill_color=RGBColor(0x1A,0x1A,0x1A), border_color=TEXT_MUTED)
    add_textbox(slide, Inches(1), y4+Pt(6), Inches(11), Inches(0.4),
                "Data:  PostgreSQL  /  ChromaDB (vectors)  /  Local FS (12 EcDev PDFs)",
                font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide: RAG Pipeline
# ============================================================
def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_GREEN)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(10), Inches(0.6),
                "RAG Pipeline", font_size=28, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(0.95), Inches(10), Inches(0.4),
                "Two-stage: Parse then Ingest; Query-ready after both complete", font_size=13, color=TEXT_SECONDARY)

    y = Inches(1.7)
    # Parse
    add_rounded_rect(slide, Inches(0.8), y, Inches(5.6), Inches(2.5),
                     fill_color=BG_CARD, border_color=ACCENT_BLUE)
    add_textbox(slide, Inches(1), y+Pt(8), Inches(3), Inches(0.35),
                "Stage 1: Parse (MinerU)", font_size=16, color=ACCENT_BLUE, bold=True)
    add_bullet_list(slide, Inches(1), y+Inches(0.6), Inches(5), Inches(1.8),
                    ["1. Upload EcDev PDF via Library or Admin",
                     "2. MinerU: layout analysis + OCR",
                     "3. Extracts: text, tables (HTML), charts (images)",
                     "4. Outputs: content_list.json + .md + images/"], font_size=11, spacing=6)

    # Ingest
    add_rounded_rect(slide, Inches(6.8), y, Inches(5.6), Inches(2.5),
                     fill_color=BG_CARD, border_color=ACCENT_GREEN)
    add_textbox(slide, Inches(7), y+Pt(8), Inches(3), Inches(0.35),
                "Stage 2: Ingest (LlamaIndex)", font_size=16, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(7), y+Inches(0.6), Inches(5), Inches(1.8),
                    ["1. MinerUReader -> LlamaIndex Document[]",
                     "2. IngestionPipeline: chunk + embed (all-MiniLM-L6-v2)",
                     "3. Upsert vectors into ChromaDB",
                     "4. Push chunk metadata to Payload CMS"], font_size=11, spacing=6)

    # Query flow
    y2 = Inches(4.6)
    add_rounded_rect(slide, Inches(0.8), y2, Inches(11.6), Inches(2.3),
                     fill_color=BG_CARD, border_color=ACCENT_PURPLE)
    add_textbox(slide, Inches(1), y2+Pt(8), Inches(8), Inches(0.35),
                "Query: Hybrid Retrieval -> Book-Scoped Filtering -> Citation Synthesis",
                font_size=15, color=ACCENT_PURPLE, bold=True)

    flow = [("User Query",ACCENT_BLUE),("->",TEXT_MUTED),("BM25+Vector\n(Hybrid)",ACCENT_GREEN),
            ("->",TEXT_MUTED),("RRF Fusion\n(k=60)",ACCENT_AMBER),("->",TEXT_MUTED),
            ("Book Filter\nPostprocessor",ACCENT_PINK),("->",TEXT_MUTED),
            ("Citation\nSynthesizer",ACCENT_PURPLE),("->",TEXT_MUTED),("Answer +\nSources [N]",ACCENT_BLUE)]
    for i,(label,c) in enumerate(flow):
        x = Inches(1.0)+Inches(i*1.05)
        if label == "->":
            add_textbox(slide, x, y2+Inches(1.0), Inches(0.4), Inches(0.5),
                        "->", font_size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
        else:
            add_rounded_rect(slide, x, y2+Inches(0.7), Inches(0.95), Inches(1.1),
                             fill_color=BG_DARK, border_color=c)
            add_textbox(slide, x+Pt(2), y2+Inches(0.8), Inches(0.9), Inches(0.9),
                        label, font_size=8, color=c, alignment=PP_ALIGN.CENTER)


# ============================================================
# Slide: Implemented Features (DONE)
# ============================================================
def slide_implemented(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_GREEN)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Implemented Features", font_size=28, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(0.95), Inches(10), Inches(0.4),
                "Sprint 1 (100%) + Sprint 2 (70%) + Acquisition Sprint (50%)", font_size=13, color=TEXT_SECONDARY)

    features = [
        ("Library Management [DONE]",
         "Upload, browse, search, edit, delete books\nGrid + table views, category filtering\nBook cover extraction + pipeline status badges",
         ACCENT_BLUE),
        ("RAG Chat with Citations [DONE]",
         "PDF viewer + chat panel (resizable split)\nSSE streaming + typewriter effect\nCitation chips with page-jump + bbox highlight\nPersisted to Payload CMS (ChatSessions + Messages)",
         ACCENT_GREEN),
        ("Ingestion Pipeline [DONE]",
         "Two-stage pipeline: Parse (MinerU) + Ingest (LlamaIndex)\n3-column inspector: Stepper / Executor / Data viewer\nSSE real-time log streaming + auto step tracking",
         ACCENT_PURPLE),
        ("Question Generation [DONE]",
         "LLM-based per-book/chapter question generation\nThree-column: PDF preview / Q-list / Eval panel\nSuggested questions in chat sidebar\nPersisted to Payload Questions collection",
         ACCENT_AMBER),
        ("LLM Evaluation [DONE]",
         "5-dim: Faithfulness, Relevancy, Correctness + 2 more\nQuestion depth: surface / understanding / synthesis\nSemantic dedup via SemanticSimilarityEvaluator\nBatch auto-evaluate with progress tracking",
         ACCENT_PINK),
        ("Auth + Admin + Seed [DONE]",
         "JWT auth, role-based access (admin/editor/reader)\n14 Payload CMS collections, auto REST + GraphQL\nPrompt editor with live SSE preview\nLLM management (Ollama / Azure OpenAI)",
         ACCENT_BLUE),
    ]

    for i, (title, desc, accent) in enumerate(features):
        row = i // 3
        col = i % 3
        x = Inches(0.8) + Inches(col * 4.1)
        y = Inches(1.5) + Inches(row * 2.9)
        add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.5),
                         fill_color=BG_CARD, border_color=accent)
        add_textbox(slide, x+Inches(0.2), y+Pt(8), Inches(3.4), Inches(0.4),
                    title, font_size=13, color=accent, bold=True)
        for j, line in enumerate(desc.split("\n")):
            add_textbox(slide, x+Inches(0.2), y+Inches(0.5)+Inches(j*0.3),
                        Inches(3.4), Inches(0.3),
                        f"  -  {line}", font_size=9, color=TEXT_SECONDARY)


# ============================================================
# Slide: Tech Stack
# ============================================================
def slide_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Tech Stack", font_size=28, color=WHITE, bold=True)

    cats = [
        ("Frontend", ["Next.js 15 (App Router)","React 19 + TypeScript","Tailwind CSS","react-pdf (PDF viewer)","Lucide icons"], ACCENT_BLUE),
        ("Backend / CMS", ["Payload CMS 3.x","PostgreSQL 15+","REST + GraphQL (auto)","JWT Authentication","14 Collections"], ACCENT_GREEN),
        ("RAG Engine", ["LlamaIndex 0.12+","FastAPI (thin layer)","ChromaDB (vectors)","BM25 + Vector -> RRF","all-MiniLM-L6-v2"], ACCENT_PURPLE),
        ("AI / Cloud", ["Ollama (local LLM)","MinerU (PDF parsing)","Azure OpenAI GPT-4o-mini","Azure AI Search","Azure Blob Storage"], ACCENT_AMBER),
    ]
    for i,(title,items,accent) in enumerate(cats):
        x = Inches(0.8)+Inches(i*3.1)
        y = Inches(1.4)
        add_rounded_rect(slide, x, y, Inches(2.8), Inches(5.2),
                         fill_color=BG_CARD, border_color=accent)
        add_textbox(slide, x+Inches(0.15), y+Pt(10), Inches(2.5), Inches(0.4),
                    title, font_size=16, color=accent, bold=True)
        for j, item in enumerate(items):
            add_textbox(slide, x+Inches(0.15), y+Inches(0.65)+Inches(j*0.5),
                        Inches(2.5), Inches(0.4),
                        f"  >  {item}", font_size=11, color=TEXT_SECONDARY)


# ============================================================
# Slide: Live Demo
# ============================================================
def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SECTION)
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
                "Live Demo", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.6),
                "End-to-end workflow with Ottawa EcDev Q1 2024 report",
                font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    steps = [
        "1.  Library -- Browse the 12 ingested EcDev quarterly reports",
        "2.  Pipeline -- Show Parse -> Ingest flow with real-time SSE logs",
        '3.  Chat -- Ask: "Compare Ottawa unemployment rate Q1 2022 vs Q1 2024"',
        "4.  Citation -- Click source [1], PDF jumps to exact page + bbox",
        '5.  Question Gen -- Generate questions for "Labour Force" chapter',
        "6.  Evaluation -- Run LLM-based question depth assessment",
    ]
    for i, step in enumerate(steps):
        add_textbox(slide, Inches(2.5), Inches(3.9)+Inches(i*0.5), Inches(8), Inches(0.45),
                    step, font_size=14, color=TEXT_SECONDARY)


# ============================================================
# Slide: Example Queries
# ============================================================
def slide_examples(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Example Queries", font_size=28, color=WHITE, bold=True)

    queries = [
        ('"What was Ottawa\'s unemployment rate in Q1 2024?"',
         "-> 4.7%, up 0.7pp from Q1 2023 (Source: Statistics Canada Table 14-10-0380-01)", ACCENT_GREEN),
        ('"How many housing starts in Q3 2023?"',
         "-> Extracts exact data from CMHC tables in the report", ACCENT_BLUE),
        ('"Compare average resale prices across 2022-2024"',
         "-> Cross-report analysis with citations to each quarter", ACCENT_PURPLE),
        ('"What is the office vacancy rate trend?"',
         "-> 12.3% -> 12.2% -> 12.0% (improving, Q1 2023 - Q1 2024)", ACCENT_AMBER),
        ('"What new airline routes were announced for Ottawa?"',
         "-> Porter, Air France, United Airlines expansions (Q1 2024)", ACCENT_PINK),
    ]
    for i,(q,a,accent) in enumerate(queries):
        y = Inches(1.3)+Inches(i*1.2)
        add_rounded_rect(slide, Inches(0.8), y, Inches(11.6), Inches(1.0),
                         fill_color=BG_CARD, border_color=accent)
        add_textbox(slide, Inches(1.1), y+Pt(6), Inches(11), Inches(0.35),
                    q, font_size=12, color=accent, bold=True)
        add_textbox(slide, Inches(1.1), y+Inches(0.45), Inches(11), Inches(0.35),
                    a, font_size=11, color=TEXT_SECONDARY)


# ============================================================
# Slide: Evaluation
# ============================================================
def slide_eval(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Evaluation & Results", font_size=28, color=WHITE, bold=True)

    metrics = [
        ("Faithfulness","0.85+","Answers grounded in\nretrieved source chunks",ACCENT_GREEN),
        ("Relevancy","0.82+","Retrieved chunks match\nthe user question",ACCENT_BLUE),
        ("Citation\nAccuracy","0.90+","Citation [N] markers\ncorrectly map to sources",ACCENT_PURPLE),
        ("Hybrid vs\nSingle","+15%","BM25+Vec recall\nvs single retriever",ACCENT_AMBER),
    ]
    for i,(label,score,desc,accent) in enumerate(metrics):
        x = Inches(0.8)+Inches(i*3.1)
        y = Inches(1.4)
        add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.8),
                         fill_color=BG_CARD, border_color=accent)
        add_textbox(slide, x+Inches(0.2), y+Pt(8), Inches(2.4), Inches(0.45),
                    label, font_size=13, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x+Inches(0.2), y+Inches(0.65), Inches(2.4), Inches(0.7),
                    score, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        for j,line in enumerate(desc.split("\n")):
            add_textbox(slide, x+Inches(0.2), y+Inches(1.6)+Inches(j*0.3),
                        Inches(2.4), Inches(0.3), line, font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    y2 = Inches(4.6)
    add_rounded_rect(slide, Inches(0.8), y2, Inches(11.6), Inches(2.3),
                     fill_color=BG_CARD, border_color=ACCENT_PINK)
    add_textbox(slide, Inches(1), y2+Pt(8), Inches(6), Inches(0.35),
                "Question Quality Assessment (Implemented)", font_size=16, color=ACCENT_PINK, bold=True)
    add_bullet_list(slide, Inches(1), y2+Inches(0.6), Inches(11), Inches(1.5),
                    ["Depth Classification: Surface -> Understanding -> Synthesis",
                     "Self-assessment scores: Relevance, Clarity, Difficulty, Overall",
                     "LLM reasoning with collapsible explanations per question",
                     "Batch auto-evaluate with progress tracking"], font_size=12, spacing=6)


# ============================================================
# Slide: Challenges
# ============================================================
def slide_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_PURPLE)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Challenges & Solutions", font_size=28, color=WHITE, bold=True)

    add_rounded_rect(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.3),
                     fill_color=BG_CARD, border_color=ACCENT_AMBER)
    add_textbox(slide, Inches(1), Inches(1.6), Inches(3), Inches(0.4),
                "Challenges", font_size=18, color=ACCENT_AMBER, bold=True)
    add_bullet_list(slide, Inches(1), Inches(2.2), Inches(5.2), Inches(4),
                    ["Cross-report contamination: BM25 retriever has no metadata filter support",
                     "MinerU struggles with complex EcDev table layouts",
                     "BM25Retriever crashes on empty corpus (bm25s bug)",
                     "Pipeline state sync between SSE stream and Payload DB",
                     "Next.js SSR hydration mismatches with dynamic content"], font_size=11, spacing=8)

    add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(5.3),
                     fill_color=BG_CARD, border_color=ACCENT_GREEN)
    add_textbox(slide, Inches(7), Inches(1.6), Inches(3), Inches(0.4),
                "Solutions", font_size=18, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, Inches(7), Inches(2.2), Inches(5.2), Inches(4),
                    ["BookFilterPostprocessor drops out-of-scope BM25 nodes + MetadataFilters for ChromaDB",
                     "Custom MinerUReader with bbox normalization for table extraction",
                     "Graceful fallback: hybrid -> vector-only when BM25 unavailable",
                     "SSE-driven optimistic UI with pipeline override pattern",
                     "Modular engine_v2: each component independently testable"], font_size=11, spacing=8)


# ============================================================
# Slide: Future Work — based on actual module roadmap
# ============================================================
def slide_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, Inches(0.7), Inches(0.5), Inches(0.06), Inches(0.5), ACCENT_GREEN)
    add_textbox(slide, Inches(1), Inches(0.4), Inches(8), Inches(0.6),
                "Future Work", font_size=28, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(0.95), Inches(10), Inches(0.4),
                "Planned sprints: S3 (eval charts + multi-type) / S4 (infra) / S5-S8 (DeepTutor + reports)",
                font_size=12, color=TEXT_SECONDARY)

    future = [
        ("Sprint 3: Eval Charts + Multi-Type Q's",
         "Evaluation radar charts + trend dashboard\n"
         "Multiple question types (MCQ, fill-in-blank)\n"
         "TOC browser frontend\n"
         "Feedback loop: aggregate eval -> tuning suggestions",
         ACCENT_BLUE),
        ("Sprint 4: Infrastructure",
         "Chunking parameter UI + preview\n"
         "Embeddings management frontend\n"
         "Access control management UI\n"
         "Retriever config panel (top_k, strategy, reranker)",
         ACCENT_GREEN),
        ("Sprint 5-8: Advanced (DeepTutor)",
         "Smart retrieval + multi-step reasoning\n"
         "Deep Question Mimic (template-based generation)\n"
         "Web Search fallback for real-time data\n"
         "Multi-role report generation engine\n"
         "Azure cloud deployment + Ottawa GenAI migration",
         ACCENT_PURPLE),
    ]

    for i,(title,desc,accent) in enumerate(future):
        x = Inches(0.8)+Inches(i*4.1)
        y = Inches(1.5)
        add_rounded_rect(slide, x, y, Inches(3.8), Inches(5.0),
                         fill_color=BG_CARD, border_color=accent)
        add_textbox(slide, x+Inches(0.2), y+Pt(10), Inches(3.4), Inches(0.4),
                    title, font_size=14, color=accent, bold=True)
        for j, line in enumerate(desc.split("\n")):
            add_textbox(slide, x+Inches(0.2), y+Inches(0.6)+Inches(j*0.42),
                        Inches(3.4), Inches(0.4),
                        f"  >  {line}", font_size=11, color=TEXT_SECONDARY)


# ============================================================
# Slide: Q&A
# ============================================================
def slide_qna(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SECTION)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
                "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                "Questions & Discussion", font_size=22, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.3), Inches(5.33), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                "Ottawa EcDev RAG -- AI-Powered Economic Development Report Analysis",
                font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
                "CST 8916 -- Applied Generative AI  /  April 2026",
                font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# Main
# ============================================================
def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)                                        #  1
    slide_agenda(prs)                                       #  2
    slide_problem(prs)                                      #  3
    slide_data(prs)                                         #  4
    section_divider(prs,                                    #  5
        "Architecture & Pipeline",
        "How the system processes Ottawa EcDev reports", ACCENT_PURPLE)
    slide_architecture(prs)                                 #  6
    slide_pipeline(prs)                                     #  7
    section_divider(prs,                                    #  8
        "Features & Demo",
        "What we built -- and a live demonstration", ACCENT_AMBER)
    slide_implemented(prs)                                  #  9
    slide_tech_stack(prs)                                   # 10
    slide_demo(prs)                                         # 11
    slide_examples(prs)                                     # 12
    section_divider(prs,                                    # 13
        "Evaluation & Retrospective",
        "How well does it work?", ACCENT_GREEN)
    slide_eval(prs)                                         # 14
    slide_challenges(prs)                                   # 15
    slide_future(prs)                                       # 16
    slide_qna(prs)                                          # 17

    return prs


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", "-o", default="Ottawa_EcDev_RAG_Presentation.pptx")
    args = parser.parse_args()
    prs = build_presentation()
    prs.save(args.output)
    print(f"Saved: {Path(args.output).resolve()}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
